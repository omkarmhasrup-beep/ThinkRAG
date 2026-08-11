import os
import csv
import docx
import PyPDF2
from pptx import Presentation
from ..config import settings
from ..models.models import File
from ..database import SessionLocal
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain_core.documents import Document
import httpx

VECTOR_STORE_DIR = "vector_stores"

def get_embeddings_model():
    return HuggingFaceEmbeddings(model_name="all-MiniLM-L6-v2")

def get_vectorstore_path(user_id: int):
    return os.path.join(VECTOR_STORE_DIR, f"user_{user_id}")

def extract_text_from_file(filepath: str, filetype: str) -> str:
    if not os.path.exists(filepath):
        return ""
        
    try:
        if filetype == "txt" or filetype == "csv":
            with open(filepath, "r", encoding="utf-8") as f:
                return f.read()
        elif filetype == "pdf":
            text = ""
            with open(filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text
        elif filetype == "docx":
            doc = docx.Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs])
        elif filetype == "pptx":
            prs = Presentation(filepath)
            text = ""
            for slide_number, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
    except Exception as e:
        print(f"Error extracting text from {filepath}: {e}")
    return ""

def process_file_and_embed(filepath: str, filetype: str, user_id: int):
    print(f"Processing and embedding {filepath} for user {user_id}...")
    text = extract_text_from_file(filepath, filetype)
    if not text.strip():
        print(f"No text extracted from {filepath}")
        return
        
    # Split text into manageable chunks
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    
    chunks = text_splitter.split_text(text)
    
    # Generate source filename for metadata
    filename = os.path.basename(filepath)
    
    documents = [Document(page_content=chunk, metadata={"source": filename}) for chunk in chunks]
    
    if not documents:
        return
        
    embeddings = get_embeddings_model()
    vs_path = get_vectorstore_path(user_id)
    
    # Make sure parent directory exists
    os.makedirs(VECTOR_STORE_DIR, exist_ok=True)
    
    if os.path.exists(vs_path):
        # Load existing FAISS index and add new docs
        vectorstore = FAISS.load_local(vs_path, embeddings, allow_dangerous_deserialization=True)
        vectorstore.add_documents(documents)
        vectorstore.save_local(vs_path)
    else:
        # Create new FAISS index
        vectorstore = FAISS.from_documents(documents, embeddings)
        vectorstore.save_local(vs_path)
        
    print(f"Successfully generated vectors for {filepath} and stored in FAISS.")

def generate_rag_response(user_id: int, question: str):
    if not settings.GROQ_API_KEY and False: # Bypassing GROQ check since we use Ollama locally now
        pass 

    vs_path = get_vectorstore_path(user_id)
    context_text = ""
    sources_used = []
    
    if os.path.exists(vs_path):
        # We have a vector store for this user!
        try:
            embeddings = get_embeddings_model()
            vectorstore = FAISS.load_local(vs_path, embeddings, allow_dangerous_deserialization=True)
            # Retrieve the top 4 most relevant chunks with scores
            docs_and_scores = vectorstore.similarity_search_with_score(question, k=4)
            
            for i, (doc, score) in enumerate(docs_and_scores):
                source = doc.metadata.get('source', 'Document')
                # FAISS returns L2 distance. 0 is perfect match.
                # Convert roughly to percentage (heuristic: max(0, 100 - score * 50))
                sim_pct = max(0, min(100, int(100 - (score * 50))))
                sources_used.append({"chunk_id": i + 1, "source": source, "content": doc.page_content, "score": sim_pct})
                context_text += f"\n--- Excerpt {i+1} from {source} ---\n{doc.page_content}\n"
        except Exception as e:
            print(f"Error during similarity search: {e}")
    
    # Fallback to naive reading if vectorstore doesn't exist
    if not context_text.strip():
        db = SessionLocal()
        try:
            files = db.query(File).filter(File.user_id == user_id).all()
            for i, file in enumerate(files):
                content = extract_text_from_file(file.filepath, file.filetype)
                if content:
                    sources_used.append({"chunk_id": i + 1, "source": file.filename, "content": content[:500], "score": 100})
                    context_text += f"\n--- Start of {file.filename} ---\n{content[:2000]}...\n" # Truncate heavily
        finally:
            db.close()

    base_prompt = (
        "You are a helpful AI assistant. "
        "CRITICAL LANGUAGE INSTRUCTION: You must match the language of the user's question. "
        "If the user asks the question in English, you MUST answer entirely in English. "
        "If the user asks the question in Marathi (even using Roman/English script like 'kay aahe' or 'mhanje kay'), you MUST answer entirely in Marathi.\n\n"
    )

    if context_text.strip():
        system_prompt = base_prompt + (
            "TASK: Answer the user's question ONLY based on the provided context below. "
            "If the answer cannot be found in the context, your answer should be exactly: 'Sorry, the provided documents do not contain the answer.' (Translate this to Marathi if the user asked in Marathi).\n\n"
            "Context:\n" + context_text
        )
    else:
        system_prompt = base_prompt + (
            "TASK: You do not have any documents to reference right now. "
            "Your answer should be exactly: 'Sorry, no relevant documents were found to answer your question.' (Translate this to Marathi if the user asked in Marathi)."
        )

    try:
        import json
        with httpx.stream(
            "POST",
            "http://localhost:11434/api/chat",
            json={
                "model": "llama3",
                "messages": [
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": question}
                ],
                "stream": True
            },
            timeout=600.0
        ) as response:
            response.raise_for_status()
            for line in response.iter_lines():
                if line:
                    try:
                        data = json.loads(line)
                        if "message" in data and "content" in data["message"]:
                            yield data["message"]["content"]
                    except:
                        pass
                        
            # Yield sources block after generation is complete
            if sources_used:
                sources_md = "\n\n```rag-context\n" + json.dumps(sources_used) + "\n```\n"
                yield sources_md
                
    except Exception as e:
        error_details = str(e)
        if hasattr(e, 'response') and e.response is not None:
            try:
                # read the stream content for error
                error_details = str(e.response.content)
            except:
                pass
        yield f"Error connecting to local Ollama (is it running?): {error_details}"

def get_rag_stats(user_id: int):
    stats = {
        "total_chunks": 0,
        "embeddings": 0,
        "vector_database": "FAISS (Local)",
        "storage_used_bytes": 0,
        "last_updated": None
    }
    
    vs_path = get_vectorstore_path(user_id)
    
    if os.path.exists(vs_path):
        try:
            embeddings_model = get_embeddings_model()
            vectorstore = FAISS.load_local(vs_path, embeddings_model, allow_dangerous_deserialization=True)
            stats["total_chunks"] = vectorstore.index.ntotal
            stats["embeddings"] = vectorstore.index.ntotal
            
            # calculate FAISS directory size
            for f in os.listdir(vs_path):
                fp = os.path.join(vs_path, f)
                if os.path.isfile(fp):
                    stats["storage_used_bytes"] += os.path.getsize(fp)
                    
            # last updated based on FAISS index file
            index_path = os.path.join(vs_path, "index.faiss")
            if os.path.exists(index_path):
                import datetime
                mtime = os.path.getmtime(index_path)
                stats["last_updated"] = datetime.datetime.fromtimestamp(mtime).isoformat()
        except Exception as e:
            print(f"Error fetching RAG stats: {e}")
            
    db = SessionLocal()
    try:
        files = db.query(File).filter(File.user_id == user_id).all()
        for file in files:
            if os.path.exists(file.filepath):
                stats["storage_used_bytes"] += os.path.getsize(file.filepath)
                # update last_updated if no FAISS index
                if not stats["last_updated"]:
                    import datetime
                    mtime = os.path.getmtime(file.filepath)
                    stats["last_updated"] = datetime.datetime.fromtimestamp(mtime).isoformat()
    finally:
        db.close()
        
    return stats

def search_knowledge_base(user_id: int, query: str, k: int = 20):
    vs_path = get_vectorstore_path(user_id)
    results = []
    
    if os.path.exists(vs_path):
        try:
            embeddings_model = get_embeddings_model()
            vectorstore = FAISS.load_local(vs_path, embeddings_model, allow_dangerous_deserialization=True)
            docs_and_scores = vectorstore.similarity_search_with_score(query, k=k)
            
            for i, (doc, score) in enumerate(docs_and_scores):
                source = doc.metadata.get('source', 'Document')
                sim_pct = max(0, min(100, int(100 - (score * 50))))
                results.append({
                    "id": i,
                    "source": source,
                    "content": doc.page_content,
                    "score": sim_pct
                })
        except Exception as e:
            print(f"Error during knowledge base search: {e}")
            
    return results
