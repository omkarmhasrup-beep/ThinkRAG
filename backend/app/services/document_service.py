import os
from langchain_core.documents import Document

def extract_text_from_file(filepath: str, filetype: str) -> str:
    if not os.path.exists(filepath):
        return ""
        
    try:
        if filetype == "txt" or filetype == "csv":
            with open(filepath, "r", encoding="utf-8") as f:
                return f.read()
        elif filetype == "pdf":
            import PyPDF2
            text = ""
            with open(filepath, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text
        elif filetype == "docx":
            import docx
            doc = docx.Document(filepath)
            return "\n".join([para.text for para in doc.paragraphs])
        elif filetype == "pptx":
            from pptx import Presentation
            prs = Presentation(filepath)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
    except Exception as e:
        print(f"Error extracting text from {filepath}: {e}")
    return ""

def chunk_text(text: str, source: str) -> list[Document]:
    if not text.strip():
        return []
        
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=1000,
        chunk_overlap=200,
        length_function=len
    )
    
    chunks = text_splitter.split_text(text)
    return [Document(page_content=chunk, metadata={"source": source}) for chunk in chunks]

